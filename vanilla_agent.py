
import os
import json
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook
import openai

# Configure OpenAI client
client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Helper Functions for File Processing
def extract_pptx_structure(file_path):
    """Extract slides from a PPTX file to create memory structure."""
    try:
        prs = Presentation(file_path)
        slides = [f"Slide {i+1}" for i in range(len(prs.slides))]
        return os.path.basename(file_path), slides
    except Exception as e:
        return os.path.basename(file_path), [f"Error: {str(e)}"]

def extract_excel_structure(file_path):
    """Extract sheet names from an Excel file to create memory structure."""
    try:
        wb = load_workbook(filename=file_path, read_only=True)
        sheet_names = wb.sheetnames
        return os.path.basename(file_path), sheet_names
    except Exception as e:
        return os.path.basename(file_path), [f"Error: {str(e)}"]

def get_slide_xml(file_path, slide_index):
    """Get XML representation of a specific slide."""
    try:
        prs = Presentation(file_path)
        if 0 <= slide_index < len(prs.slides):
            # Create a simplified XML representation of the slide
            slide = prs.slides[slide_index]
            xml_representation = "<slide>\n"
            
            # Add shapes
            xml_representation += "  <shapes>\n"
            for i, shape in enumerate(slide.shapes):
                shape_type = type(shape).__name__
                xml_representation += f"    <shape id='{i}' type='{shape_type}'>\n"
                
                # Add text if present
                if hasattr(shape, "text_frame") and shape.text_frame:
                    xml_representation += "      <text_frame>\n"
                    for paragraph in shape.text_frame.paragraphs:
                        xml_representation += f"        <paragraph>{paragraph.text}</paragraph>\n"
                    xml_representation += "      </text_frame>\n"
                
                # Add table if present
                if hasattr(shape, "table") and shape.table:
                    xml_representation += "      <table>\n"
                    for row in shape.table.rows:
                        xml_representation += "        <row>\n"
                        for cell in row.cells:
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            xml_representation += f"          <cell>{cell_text}</cell>\n"
                        xml_representation += "        </row>\n"
                    xml_representation += "      </table>\n"
                
                xml_representation += "    </shape>\n"
            xml_representation += "  </shapes>\n"
            
            xml_representation += "</slide>"
            return xml_representation
        else:
            return f"Error: Slide index {slide_index} out of range."
    except Exception as e:
        return f"Error: {str(e)}"

def get_excel_table(file_path, sheet_name):
    """Get markdown table representation of an Excel sheet."""
    try:
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        return df.to_markdown(index=False)
    except Exception as e:
        return f"Error: {str(e)}"

def modify_slide(file_path, slide_index, code):
    """Modify a slide using Python code."""
    try:
        # Create a local variable to hold the presentation
        prs = Presentation(file_path)
        if 0 <= slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            
            # Execute the code in a context with access to slide
            local_vars = {"slide": slide, "prs": prs}
            exec(code, {}, local_vars)
            
            # Save the modified presentation
            prs.save(file_path)
            
            # Return the updated XML
            return get_slide_xml(file_path, slide_index)
        else:
            return f"Error: Slide index {slide_index} out of range."
    except Exception as e:
        return f"Error: {str(e)}\n\nCode attempted to execute:\n{code}"

def modify_excel(file_path, sheet_name, code):
    """Modify an Excel sheet using Python code."""
    try:
        # Read the sheet into a DataFrame
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        
        # Execute the code with access to the DataFrame
        local_vars = {"df": df}
        exec(code, {}, local_vars)
        
        # Get the updated DataFrame
        updated_df = local_vars.get("df", df)
        
        # Write the updated DataFrame back to the Excel file
        with pd.ExcelWriter(file_path, engine='openpyxl', mode='a', if_sheet_exists='replace') as writer:
            updated_df.to_excel(writer, sheet_name=sheet_name, index=False)
        
        # Return the updated table
        return updated_df.to_markdown(index=False)
    except Exception as e:
        return f"Error: {str(e)}\n\nCode attempted to execute:\n{code}"

# LLM Tools Definition
def define_tools():
    """Define tools for the LLM to interact with files."""
    return [
        {
            "type": "function",
            "function": {
                "name": "get_slide",
                "description": "Get the XML representation of a slide from a PowerPoint file",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Path to the PowerPoint file"
                        },
                        "slide_index": {
                            "type": "integer",
                            "description": "Zero-based index of the slide to retrieve"
                        }
                    },
                    "required": ["file_path", "slide_index"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "get_excel_data",
                "description": "Get the data from an Excel sheet as a markdown table",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Path to the Excel file"
                        },
                        "sheet_name": {
                            "type": "string",
                            "description": "Name of the sheet to retrieve"
                        }
                    },
                    "required": ["file_path", "sheet_name"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "modify_slide",
                "description": "Modify a slide using Python code",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Path to the PowerPoint file"
                        },
                        "slide_index": {
                            "type": "integer",
                            "description": "Zero-based index of the slide to modify"
                        },
                        "code": {
                            "type": "string",
                            "description": "Python code to execute to modify the slide (has access to 'slide' object from python-pptx)"
                        }
                    },
                    "required": ["file_path", "slide_index", "code"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "modify_excel",
                "description": "Modify an Excel sheet using Python code",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Path to the Excel file"
                        },
                        "sheet_name": {
                            "type": "string",
                            "description": "Name of the sheet to modify"
                        },
                        "code": {
                            "type": "string",
                            "description": "Python code to execute to modify the sheet (has access to 'df' DataFrame from pandas)"
                        }
                    },
                    "required": ["file_path", "sheet_name", "code"]
                }
            }
        }
    ]

# Execute LLM Tool
def execute_tool(tool_name, tool_args):
    """Execute the tool called by the LLM."""
    if tool_name == "get_slide":
        return get_slide_xml(tool_args["file_path"], tool_args["slide_index"])
    elif tool_name == "get_excel_data":
        return get_excel_table(tool_args["file_path"], tool_args["sheet_name"])
    elif tool_name == "modify_slide":
        return modify_slide(tool_args["file_path"], tool_args["slide_index"], tool_args["code"])
    elif tool_name == "modify_excel":
        return modify_excel(tool_args["file_path"], tool_args["sheet_name"], tool_args["code"])
    else:
        return f"Unknown tool: {tool_name}"

# Create Memory Snapshot
def create_memory_snapshot(pptx_files, excel_files):
    """Create a memory snapshot structure for the LLM."""
    memory = {"Memory": {}}
    
    for file_path in pptx_files:
        deck_name, slides = extract_pptx_structure(file_path)
        memory["Memory"][deck_name] = slides
    
    for file_path in excel_files:
        workbook_name, sheets = extract_excel_structure(file_path)
        memory["Memory"][workbook_name] = sheets
    
    return memory

# Create a mapping between file names and their full paths
def create_file_path_mapping(pptx_files, excel_files):
    """Create a mapping between file names and their full paths."""
    mapping = {}
    for file_path in pptx_files + excel_files:
        mapping[os.path.basename(file_path)] = file_path
    return mapping

# Main LLM lifecycle function
def ai_ppt_agent(user_query, pptx_files=None, excel_files=None, max_iterations=10):
    """
    Main function for the AI PPT agent.
    
    Args:
        user_query (str): The user's query
        pptx_files (list): List of PowerPoint file paths
        excel_files (list): List of Excel file paths
        max_iterations (int): Maximum number of LLM iterations
    """
    if pptx_files is None:
        pptx_files = []
    if excel_files is None:
        excel_files = []
    
    # Create memory snapshot
    memory = create_memory_snapshot(pptx_files, excel_files)
    memory_str = json.dumps(memory, indent=4)
    
    # Create file path mapping
    file_path_mapping = create_file_path_mapping(pptx_files, excel_files)
    
    # Define tools
    tools = define_tools()
    
    # Initialize conversation
    messages = [
        {
            "role": "system",
            "content": f"""You are an AI PowerPoint and Excel agent. You can view and modify PowerPoint slides and Excel sheets.
            
The memory snapshot of available files is:
{memory_str}

File paths mapping:
{json.dumps(file_path_mapping, indent=2)}

You have access to the following tools:
1. get_slide - Get the XML representation of a slide
2. get_excel_data - Get data from an Excel sheet as a markdown table
3. modify_slide - Modify a slide using Python code
4. modify_excel - Modify an Excel sheet using Python code

When modifying slides, you have access to a 'slide' object from the python-pptx library.
When modifying Excel, you have access to a 'df' DataFrame object from pandas.

Always plan your approach before making changes. First examine the files to understand their structure,
then make targeted modifications based on the user's request.
"""
        },
        {
            "role": "user",
            "content": user_query
        }
    ]
    
    iteration = 0
    while iteration < max_iterations:
        print(f"Iteration {iteration + 1}/{max_iterations}")
        
        # Call the LLM
        response = client.chat.completions.create(
            model="o3-mini",
            messages=messages,
            tools=tools,
            tool_choice="auto"
        )
        
        assistant_message = response.choices[0].message
        
        # Add the assistant's message to the conversation
        messages.append({
            "role": "assistant",
            "content": assistant_message.content or "",
            "tool_calls": assistant_message.tool_calls if hasattr(assistant_message, 'tool_calls') else None
        })
        
        # Check if the assistant wants to use a tool
        if hasattr(assistant_message, 'tool_calls') and assistant_message.tool_calls:
            for tool_call in assistant_message.tool_calls:
                tool_name = tool_call.function.name
                tool_args = json.loads(tool_call.function.arguments)
                
                print(f"Using tool: {tool_name}")
                
                # Execute the tool
                tool_response = execute_tool(tool_name, tool_args)
                
                # Update the memory if modification tools were used
                if tool_name in ["modify_slide", "modify_excel"]:
                    memory = create_memory_snapshot(pptx_files, excel_files)
                    memory_str = json.dumps(memory, indent=4)
                    
                    # Update the system message with the new memory
                    messages[0]["content"] = f"""You are an AI PowerPoint and Excel agent. You can view and modify PowerPoint slides and Excel sheets.
                    
The memory snapshot of available files is:
{memory_str}

File paths mapping:
{json.dumps(file_path_mapping, indent=2)}

You have access to the following tools:
1. get_slide - Get the XML representation of a slide
2. get_excel_data - Get data from an Excel sheet as a markdown table
3. modify_slide - Modify a slide using Python code
4. modify_excel - Modify an Excel sheet using Python code

When modifying slides, you have access to a 'slide' object from the python-pptx library.
When modifying Excel, you have access to a 'df' DataFrame object from pandas.

Always plan your approach before making changes. First examine the files to understand their structure,
then make targeted modifications based on the user's request.
"""
                
                # Add the tool response to messages
                messages.append({
                    "role": "tool",
                    "tool_call_id": tool_call.id,
                    "name": tool_name,
                    "content": str(tool_response)
                })
            
            # Continue to next iteration to process tool responses
            continue
        
        # If the assistant doesn't call a tool, show response and get user input
        print("Assistant:", assistant_message.content)
        user_input = input("You (type 'exit' to end): ")
        
        if user_input.lower() == 'exit':
            break
        
        # Add the user's response to messages
        messages.append({
            "role": "user",
            "content": user_input
        })
        
        iteration += 1
    
    return messages

# Example usage
if __name__ == "__main__":
    user_query = input("What would you like to do with your PowerPoint or Excel files? ")
    pptx_files = input("Enter PowerPoint file paths (comma-separated, or press Enter if none): ").split(',')
    excel_files = input("Enter Excel file paths (comma-separated, or press Enter if none): ").split(',')
    
    pptx_files = [f.strip() for f in pptx_files if f.strip()]
    excel_files = [f.strip() for f in excel_files if f.strip()]
    
    conversation = ai_ppt_agent(user_query, pptx_files, excel_files)
    
    # Save the conversation if needed
    # with open("conversation_log.json", "w") as f:
    #     json.dump(conversation, f, indent=4)
