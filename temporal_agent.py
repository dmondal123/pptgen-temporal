import asyncio
from temporalio import workflow, activity
from temporalio.client import Client
from temporalio.worker import Worker
from datetime import timedelta
from typing import List, Dict, Any, Optional


# Helper Functions for File Processing
@activity.defn
async def extract_pptx_structure(file_path: str) -> tuple:
    """Extract slides from a PPTX file to create memory structure."""
    import os
    from pptx import Presentation
    
    try:
        prs = Presentation(file_path)
        slides = [f"Slide {i+1}" for i in range(len(prs.slides))]
        return os.path.basename(file_path), slides
    except Exception as e:
        return os.path.basename(file_path), [f"Error: {str(e)}"]

@activity.defn
async def extract_excel_structure(file_path: str) -> tuple:
    """Extract sheet names from an Excel file to create memory structure."""
    import os
    from openpyxl import load_workbook
    
    try:
        wb = load_workbook(filename=file_path, read_only=True)
        sheet_names = wb.sheetnames
        return os.path.basename(file_path), sheet_names
    except Exception as e:
        return os.path.basename(file_path), [f"Error: {str(e)}"]

@activity.defn
async def get_slide_xml(file_path: str, slide_index: int) -> str:
    """Get XML representation of a specific slide."""
    from pptx import Presentation
    
    try:
        prs = Presentation(file_path)
        if 0 <= slide_index < len(prs.slides):
            # Create a simplified XML representation of the slide
            slide = prs.slides[slide_index]
            xml_representation = "<slide>\n"
            
            # Add shapes
            xml_representation += "  <shapes>\n"
            for i, shape in enumerate(slide.shapes):
                shape_type = type(shape).__name__
                xml_representation += f"    <shape id='{i}' type='{shape_type}'>\n"
                
                # Add text if present
                if hasattr(shape, "text_frame") and shape.text_frame:
                    xml_representation += "      <text_frame>\n"
                    for paragraph in shape.text_frame.paragraphs:
                        xml_representation += f"        <paragraph>{paragraph.text}</paragraph>\n"
                    xml_representation += "      </text_frame>\n"
                
                # Add table if present
                if hasattr(shape, "table") and shape.table:
                    xml_representation += "      <table>\n"
                    for row in shape.table.rows:
                        xml_representation += "        <row>\n"
                        for cell in row.cells:
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            xml_representation += f"          <cell>{cell_text}</cell>\n"
                        xml_representation += "        </row>\n"
                    xml_representation += "      </table>\n"
                
                xml_representation += "    </shape>\n"
            xml_representation += "  </shapes>\n"
            
            xml_representation += "</slide>"
            return xml_representation
        else:
            return f"Error: Slide index {slide_index} out of range."
    except Exception as e:
        return f"Error: {str(e)}"

@activity.defn
async def get_excel_table(file_path: str, sheet_name: str) -> str:
    """Get markdown table representation of an Excel sheet."""
    import pandas as pd
    
    try:
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        return df.to_markdown(index=False)
    except Exception as e:
        return f"Error: {str(e)}"

@activity.defn
async def modify_slide(file_path: str, slide_index: int, code: str) -> str:
    """Modify a slide using Python code."""
    from pptx import Presentation
    
    try:
        # Create a local variable to hold the presentation
        prs = Presentation(file_path)
        if 0 <= slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            
            # Execute the code in a context with access to slide
            local_vars = {"slide": slide, "prs": prs}
            exec(code, {}, local_vars)
            
            # Save the modified presentation
            prs.save(file_path)
            
            # Return the updated XML - inline the XML generation instead of calling get_slide_xml
            slide = prs.slides[slide_index]
            xml_representation = "<slide>\n"
            
            # Add shapes
            xml_representation += "  <shapes>\n"
            for i, shape in enumerate(slide.shapes):
                shape_type = type(shape).__name__
                xml_representation += f"    <shape id='{i}' type='{shape_type}'>\n"
                
                # Add text if present
                if hasattr(shape, "text_frame") and shape.text_frame:
                    xml_representation += "      <text_frame>\n"
                    for paragraph in shape.text_frame.paragraphs:
                        xml_representation += f"        <paragraph>{paragraph.text}</paragraph>\n"
                    xml_representation += "      </text_frame>\n"
                
                # Add table if present
                if hasattr(shape, "table") and shape.table:
                    xml_representation += "      <table>\n"
                    for row in shape.table.rows:
                        xml_representation += "        <row>\n"
                        for cell in row.cells:
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            xml_representation += f"          <cell>{cell_text}</cell>\n"
                        xml_representation += "        </row>\n"
                    xml_representation += "      </table>\n"
                
                xml_representation += "    </shape>\n"
            xml_representation += "  </shapes>\n"
            
            xml_representation += "</slide>"
            return xml_representation
        else:
            return f"Error: Slide index {slide_index} out of range."
    except Exception as e:
        return f"Error: {str(e)}\n\nCode attempted to execute:\n{code}"

@activity.defn
async def modify_excel(file_path: str, sheet_name: str, code: str) -> str:
    """Modify an Excel sheet using Python code."""
    import pandas as pd
    
    try:
        # Read the sheet into a DataFrame
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        
        # Execute the code with access to the DataFrame
        local_vars = {"df": df}
        exec(code, {}, local_vars)
        
        # Get the updated DataFrame
        updated_df = local_vars.get("df", df)
        
        # Write the updated DataFrame back to the Excel file
        with pd.ExcelWriter(file_path, engine='openpyxl', mode='a', if_sheet_exists='replace') as writer:
            updated_df.to_excel(writer, sheet_name=sheet_name, index=False)
        
        # Return the updated table
        return updated_df.to_markdown(index=False)
    except Exception as e:
        return f"Error: {str(e)}\n\nCode attempted to execute:\n{code}"

# LLM Tools Definition
def define_tools():
    """Define tools for the LLM to interact with files."""
    return [
        {
            "type": "function",
            "function": {
                "name": "get_slide",
                "description": "Get the XML representation of a slide from a PowerPoint file",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Path to the PowerPoint file"
                        },
                        "slide_index": {
                            "type": "integer",
                            "description": "Zero-based index of the slide to retrieve"
                        }
                    },
                    "required": ["file_path", "slide_index"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "get_excel_data",
                "description": "Get the data from an Excel sheet as a markdown table",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Path to the Excel file"
                        },
                        "sheet_name": {
                            "type": "string",
                            "description": "Name of the sheet to retrieve"
                        }
                    },
                    "required": ["file_path", "sheet_name"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "modify_slide",
                "description": "Modify a slide using Python code",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Path to the PowerPoint file"
                        },
                        "slide_index": {
                            "type": "integer",
                            "description": "Zero-based index of the slide to modify"
                        },
                        "code": {
                            "type": "string",
                            "description": "Python code to execute to modify the slide (has access to 'slide' object from python-pptx)"
                        }
                    },
                    "required": ["file_path", "slide_index", "code"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "modify_excel",
                "description": "Modify an Excel sheet using Python code",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Path to the Excel file"
                        },
                        "sheet_name": {
                            "type": "string",
                            "description": "Name of the sheet to modify"
                        },
                        "code": {
                            "type": "string",
                            "description": "Python code to execute to modify the sheet (has access to 'df' DataFrame from pandas)"
                        }
                    },
                    "required": ["file_path", "sheet_name", "code"]
                }
            }
        }
    ]

@activity.defn
async def create_memory_snapshot(pptx_files: List[str], excel_files: List[str]) -> Dict:
    """Create a memory snapshot structure for the LLM."""
    import os
    from pptx import Presentation
    from openpyxl import load_workbook
    
    memory = {"Memory": {}}
    
    # Instead of calling other activities, inline the structure extraction
    for file_path in pptx_files:
        try:
            prs = Presentation(file_path)
            slides = [f"Slide {i+1}" for i in range(len(prs.slides))]
            deck_name = os.path.basename(file_path)
            memory["Memory"][deck_name] = slides
        except Exception as e:
            deck_name = os.path.basename(file_path)
            memory["Memory"][deck_name] = [f"Error: {str(e)}"]
    
    for file_path in excel_files:
        try:
            wb = load_workbook(filename=file_path, read_only=True)
            sheet_names = wb.sheetnames
            workbook_name = os.path.basename(file_path)
            memory["Memory"][workbook_name] = sheet_names
        except Exception as e:
            workbook_name = os.path.basename(file_path)
            memory["Memory"][workbook_name] = [f"Error: {str(e)}"]
    
    return memory

@activity.defn
async def create_file_path_mapping(pptx_files: List[str], excel_files: List[str]) -> Dict:
    """Create a mapping between file names and their full paths."""
    import os
    
    mapping = {}
    for file_path in pptx_files + excel_files:
        mapping[os.path.basename(file_path)] = file_path
    return mapping

@activity.defn
async def call_llm(messages: List[Dict[str, Any]], tools: List[Dict]) -> Dict:
    """Call the LLM with the given messages and tools."""
    import os
    import openai
    
    # Configure OpenAI client inside the activity
    client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    
    response = client.chat.completions.create(
        model="o3-mini",
        messages=messages,
        tools=tools,
        tool_choice="auto"
    )
    
    assistant_message = response.choices[0].message
    
    return {
        "content": assistant_message.content or "",
        "tool_calls": assistant_message.tool_calls if hasattr(assistant_message, 'tool_calls') else None
    }

@activity.defn
async def execute_tool(tool_name: str, tool_args: Dict) -> str:
    """Execute the tool called by the LLM."""
    if tool_name == "get_slide":
        # Call the get_slide_xml activity via workflow, not directly
        return await get_slide_xml(tool_args["file_path"], tool_args["slide_index"])
    elif tool_name == "get_excel_data":
        return await get_excel_table(tool_args["file_path"], tool_args["sheet_name"])
    elif tool_name == "modify_slide":
        return await modify_slide(tool_args["file_path"], tool_args["slide_index"], tool_args["code"])
    elif tool_name == "modify_excel":
        return await modify_excel(tool_args["file_path"], tool_args["sheet_name"], tool_args["code"])
    else:
        return f"Unknown tool: {tool_name}"

# Workflow Definition
@workflow.defn
class PPTAgentWorkflow:
    def __init__(self):
        self.messages = []
        self.pptx_files = []
        self.excel_files = []
        self.file_path_mapping = {}
        self.memory = {}
        self.tools = define_tools()
        self.user_input_received = False
        self.user_query = ""
        
    @workflow.run
    async def run(self) -> List[Dict[str, Any]]:
        """Main workflow that orchestrates the agent."""
        import json
        
        # Initialize with system message
        self.messages = [
            {
                "role": "system",
                "content": "You are an AI PowerPoint and Excel agent. You can view and modify PowerPoint slides and Excel sheets."
            }
        ]
        
        # Main loop
        while True:
            # Wait for user input
            await workflow.wait_condition(lambda: bool(self.user_input_received))
            self.user_input_received = False
            
            # Update memory snapshot
            self.memory = await workflow.execute_activity(
                create_memory_snapshot,
                args=[self.pptx_files, self.excel_files],
                start_to_close_timeout=timedelta(seconds=30)
            )
            
            # Update file path mapping
            self.file_path_mapping = await workflow.execute_activity(
                create_file_path_mapping,
                args=[self.pptx_files, self.excel_files],
                start_to_close_timeout=timedelta(seconds=30)
            )
            
            # Update system message with memory
            memory_str = json.dumps(self.memory, indent=4)
            file_mapping_str = json.dumps(self.file_path_mapping, indent=2)
            
            self.messages[0]["content"] = f"""You are an AI PowerPoint and Excel agent. You can view and modify PowerPoint slides and Excel sheets.
            
The memory snapshot of available files is:
{memory_str}

File paths mapping:
{file_mapping_str}

You have access to the following tools:
1. get_slide - Get the XML representation of a slide. Use this tool to examine slide structure/content before proceeding with any slide modifications.
2. get_excel_data - Get data from an Excel sheet as a markdown table. Use this tool to analyze data structure before proceeding with any excel modifications.
3. modify_slide - Use this tool to modify any Slide using Python code.
4. modify_excel - Use this tool to modify any Excel sheet using Python code.

When modifying slides, you have access to a 'slide' object from the python-pptx library.
When modifying Excel, you have access to a 'df' DataFrame object from pandas.

Always plan your approach before making changes. First examine the files to understand their structure,
then make targeted modifications based on the user's request.
"""
            
            # Add user query to messages
            self.messages.append({
                "role": "user",
                "content": self.user_query
            })
            
            # Process message chain until LLM stops calling tools
            while True:
                # Call LLM
                assistant_message = await workflow.execute_activity(
                    call_llm,
                    args=[self.messages, self.tools],
                    start_to_close_timeout=timedelta(minutes=2)
                )
                
                # Add assistant message to conversation
                self.messages.append({
                    "role": "assistant",
                    "content": assistant_message["content"],
                    "tool_calls": assistant_message["tool_calls"]
                })
                
                # Check if tools need to be called
                if assistant_message["tool_calls"]:
                    for tool_call in assistant_message["tool_calls"]:
                        tool_name = tool_call["function"]["name"]
                        tool_args = json.loads(tool_call["function"]["arguments"])
                        
                        # Execute tool
                        tool_response = await workflow.execute_activity(
                            execute_tool,
                            args=[tool_name, tool_args],
                            start_to_close_timeout=timedelta(minutes=1)
                        )
                        
                        # Update memory if modification tools were used
                        if tool_name in ["modify_slide", "modify_excel"]:
                            self.memory = await workflow.execute_activity(
                                create_memory_snapshot,
                                args=[self.pptx_files, self.excel_files],
                                start_to_close_timeout=timedelta(seconds=30)
                            )
                            
                            # Update system message with new memory
                            memory_str = json.dumps(self.memory, indent=4)
                            file_mapping_str = json.dumps(self.file_path_mapping, indent=2)
                            
                            self.messages[0]["content"] = f"""You are an AI PowerPoint and Excel agent. You can view and modify PowerPoint slides and Excel sheets.
                            
The memory snapshot of available files is:
{memory_str}

File paths mapping:
{file_mapping_str}

You have access to the following tools:
1. get_slide - Get the XML representation of a slide
2. get_excel_data - Get data from an Excel sheet as a markdown table
3. modify_slide - Modify a slide using Python code
4. modify_excel - Modify an Excel sheet using Python code

When modifying slides, you have access to a 'slide' object from the python-pptx library.
When modifying Excel, you have access to a 'df' DataFrame object from pandas.

Always plan your approach before making changes. First examine the files to understand their structure,
then make targeted modifications based on the user's request.
"""
                        
                        # Add tool response to messages
                        self.messages.append({
                            "role": "tool",
                            "tool_call_id": tool_call["id"],
                            "name": tool_name,
                            "content": str(tool_response)
                        })
                        
                        # Continue to next iteration to let LLM process tool responses
                        continue
                    
                    # Continue to next iteration to let LLM process tool responses
                    continue
                
                # If no tools were called, break the loop and wait for user input
                break
            
            # Wait for next user input
    
    @workflow.signal
    async def user_input(self, input_data: Dict[str, Any]):
        """Signal handler for user input - takes a single dictionary containing all data."""
        self.user_query = input_data.get("query", "")
        self.pptx_files = input_data.get("pptx_files", [])
        self.excel_files = input_data.get("excel_files", [])
        self.user_input_received = True
    
    @workflow.query
    def get_conversation_history(self) -> List[Dict[str, Any]]:
        """Query method to get the current conversation history."""
        return self.messages

# Worker setup and registration
async def run_worker():
    # Initialize client
    client = await Client.connect("localhost:7233")
    
    # Register workflow and activities with worker
    worker = Worker(
        client,
        task_queue="ppt-agent-task-queue",
        workflows=[PPTAgentWorkflow],
        activities=[
            extract_pptx_structure, 
            extract_excel_structure, 
            get_slide_xml, 
            get_excel_table, 
            modify_slide, 
            modify_excel,
            create_memory_snapshot,
            create_file_path_mapping,
            call_llm,
            execute_tool
        ]
    )
    
    print("Worker started. Press Ctrl+C to exit.")
    # Start worker
    await worker.run()

# Sample code for testing
async def run_test_client():
    # Connect to Temporal server
    client = await Client.connect("localhost:7233")
    
    # Start the workflow with a specific ID
    workflow_id = "test-ppt-agent-workflow"
    handle = await client.start_workflow(
        PPTAgentWorkflow.run,
        id=workflow_id,
        task_queue="ppt-agent-task-queue"
    )
    
    print(f"Started workflow with ID: {workflow_id}")
    
    # Send a signal with user query and files - Using a simple dictionary
    user_query = "what are the projects mentioned in the slide?"
    pptx_files = ["Anil Turaga - About me.pptx"]
    excel_files = ["Internship Projects Tracking.xlsx"]
    excel_files = []
    
    # Package all data in a simple dictionary
    input_data = {
        "query": user_query,
        "pptx_files": pptx_files,
        "excel_files": excel_files
    }
    
    # Send signal with a single argument
    await handle.signal(PPTAgentWorkflow.user_input, input_data)
    
    print(f"Sent signal with query: {user_query}")
    
    # Poll for conversation history
    for _ in range(5):  # Poll a few times to see updates
        await asyncio.sleep(3)  # Wait 3 seconds between polls
        history = await handle.query(PPTAgentWorkflow.get_conversation_history)
        print("\nCurrent conversation history:")
        for msg in history:
            role = msg.get("role", "")
            content = msg.get("content", "")
            print(f"{role.upper()}: {content}...")  # Print first 100 chars
    
    print("\nTest complete.")

if __name__ == "__main__":
    import sys
    
    # Check command line arguments
    if len(sys.argv) > 1 and sys.argv[1] == "test":
        print("Running in test client mode...")
        asyncio.run(run_test_client())
    else:
        print("Running in worker mode...")
        asyncio.run(run_worker()) 
