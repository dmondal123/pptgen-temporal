import asyncio
from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional
from temporalio import activity
from pptx import Presentation
from openpyxl import load_workbook
import pandas as pd
import os
import openai
import json

# Data Classes for Activity Parameters
@dataclass
class SlideParams:
    file_path: str
    slide_index: int

@dataclass
class ExcelParams:
    file_path: str
    sheet_name: str

@dataclass
class ModifySlideParams:
    file_path: str
    slide_index: int
    code: str

@dataclass
class ModifyExcelParams:
    file_path: str
    sheet_name: str
    code: str

@dataclass
class MemorySnapshotParams:
    pptx_files: List[str]
    excel_files: List[str]

@dataclass
class LLMParams:
    messages: List[Dict[str, Any]]
    tools: List[Dict]

@dataclass
class ToolExecutionParams:
    tool_name: str
    tool_args: Dict

# File Processing Activities
@activity.defn
async def extract_pptx_structure(file_path: str) -> tuple:
    """Extract slides from a PPTX file to create memory structure."""
    try:
        prs = Presentation(file_path)
        slides = [f"Slide {i+1}" for i in range(len(prs.slides))]
        return os.path.basename(file_path), slides
    except Exception as e:
        return os.path.basename(file_path), [f"Error: {str(e)}"]

@activity.defn
async def extract_excel_structure(file_path: str) -> tuple:
    """Extract sheet names from an Excel file to create memory structure."""
    try:
        wb = load_workbook(filename=file_path, read_only=True)
        sheet_names = wb.sheetnames
        return os.path.basename(file_path), sheet_names
    except Exception as e:
        return os.path.basename(file_path), [f"Error: {str(e)}"]

@activity.defn
async def get_slide_xml(params: SlideParams) -> str:
    """Get XML representation of a specific slide."""
    try:
        prs = Presentation(params.file_path)
        if 0 <= params.slide_index < len(prs.slides):
            slide = prs.slides[params.slide_index]
            return create_slide_xml(slide)
        else:
            return f"Error: Slide index {params.slide_index} out of range."
    except Exception as e:
        return f"Error: {str(e)}"

@activity.defn
async def get_excel_table(params: ExcelParams) -> str:
    """Get markdown table representation of an Excel sheet."""
    try:
        df = pd.read_excel(params.file_path, sheet_name=params.sheet_name)
        return df.to_markdown(index=False)
    except Exception as e:
        return f"Error: {str(e)}"

@activity.defn
async def modify_slide(params: ModifySlideParams) -> str:
    """Modify a slide using Python code."""
    try:
        prs = Presentation(params.file_path)
        if 0 <= params.slide_index < len(prs.slides):
            slide = prs.slides[params.slide_index]
            local_vars = {"slide": slide, "prs": prs}
            exec(params.code, {}, local_vars)
            prs.save(params.file_path)
            return create_slide_xml(slide)
        else:
            return f"Error: Slide index {params.slide_index} out of range."
    except Exception as e:
        return f"Error: {str(e)}\n\nCode attempted to execute:\n{params.code}"

@activity.defn
async def modify_excel(params: ModifyExcelParams) -> str:
    """Modify an Excel sheet using Python code."""
    try:
        df = pd.read_excel(params.file_path, sheet_name=params.sheet_name)
        local_vars = {"df": df}
        exec(params.code, {}, local_vars)
        updated_df = local_vars.get("df", df)
        
        with pd.ExcelWriter(params.file_path, engine='openpyxl', mode='a', if_sheet_exists='replace') as writer:
            updated_df.to_excel(writer, sheet_name=params.sheet_name, index=False)
        
        return updated_df.to_markdown(index=False)
    except Exception as e:
        return f"Error: {str(e)}\n\nCode attempted to execute:\n{params.code}"

# Memory and State Management Activities
@activity.defn
async def create_memory_snapshot(params: MemorySnapshotParams) -> Dict:
    """Create a memory snapshot structure for the LLM."""
    memory = {"Memory": {}}
    
    for file_path in params.pptx_files:
        try:
            prs = Presentation(file_path)
            slides = [f"Slide {i+1}" for i in range(len(prs.slides))]
            deck_name = os.path.basename(file_path)
            memory["Memory"][deck_name] = slides
        except Exception as e:
            deck_name = os.path.basename(file_path)
            memory["Memory"][deck_name] = [f"Error: {str(e)}"]
    
    for file_path in params.excel_files:
        try:
            wb = load_workbook(filename=file_path, read_only=True)
            sheet_names = wb.sheetnames
            workbook_name = os.path.basename(file_path)
            memory["Memory"][workbook_name] = sheet_names
        except Exception as e:
            workbook_name = os.path.basename(file_path)
            memory["Memory"][workbook_name] = [f"Error: {str(e)}"]
    
    return memory

# LLM Interaction Activities
@activity.defn
async def call_llm(params: LLMParams) -> Dict:
    """Call the LLM with the given messages and tools."""
    client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    
    response = client.chat.completions.create(
        model="o3-mini",
        messages=params.messages,
        tools=params.tools,
        tool_choice="auto"
    )
    
    assistant_message = response.choices[0].message
    
    return {
        "content": assistant_message.content or "",
        "tool_calls": assistant_message.tool_calls if hasattr(assistant_message, 'tool_calls') else None
    }

@activity.defn
async def execute_tool(params: ToolExecutionParams) -> str:
    """Execute the tool called by the LLM."""
    if params.tool_name == "get_slide":
        return await get_slide_xml(SlideParams(
            file_path=params.tool_args["file_path"],
            slide_index=params.tool_args["slide_index"]
        ))
    elif params.tool_name == "get_excel_data":
        return await get_excel_table(ExcelParams(
            file_path=params.tool_args["file_path"],
            sheet_name=params.tool_args["sheet_name"]
        ))
    elif params.tool_name == "modify_slide":
        return await modify_slide(ModifySlideParams(
            file_path=params.tool_args["file_path"],
            slide_index=params.tool_args["slide_index"],
            code=params.tool_args["code"]
        ))
    elif params.tool_name == "modify_excel":
        return await modify_excel(ModifyExcelParams(
            file_path=params.tool_args["file_path"],
            sheet_name=params.tool_args["sheet_name"],
            code=params.tool_args["code"]
        ))
    else:
        return f"Unknown tool: {params.tool_name}"

# Helper Functions
def create_slide_xml(slide) -> str:
    """Helper function to create XML representation of a slide."""
    xml_representation = "<slide>\n"
    xml_representation += "  <shapes>\n"
    for i, shape in enumerate(slide.shapes):
        shape_type = type(shape).__name__
        xml_representation += f"    <shape id='{i}' type='{shape_type}'>\n"
        
        if hasattr(shape, "text_frame") and shape.text_frame:
            xml_representation += "      <text_frame>\n"
            for paragraph in shape.text_frame.paragraphs:
                xml_representation += f"        <paragraph>{paragraph.text}</paragraph>\n"
            xml_representation += "      </text_frame>\n"
        
        if hasattr(shape, "table") and shape.table:
            xml_representation += "      <table>\n"
            for row in shape.table.rows:
                xml_representation += "        <row>\n"
                for cell in row.cells:
                    cell_text = cell.text_frame.text if cell.text_frame else ""
                    xml_representation += f"          <cell>{cell_text}</cell>\n"
                xml_representation += "        </row>\n"
            xml_representation += "      </table>\n"
        
        xml_representation += "    </shape>\n"
    xml_representation += "  </shapes>\n"
    xml_representation += "</slide>"
    return xml_representation
